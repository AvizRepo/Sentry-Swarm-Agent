import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # 16:9 widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Premium Dark Mode Theme Colors
    BG_COLOR = RGBColor(10, 15, 28)          # Deep slate dark background
    CARD_BG = RGBColor(21, 29, 48)           # Lighter card container background
    BORDER_COLOR = RGBColor(39, 52, 80)      # Border separator slate
    
    TEXT_WHITE = RGBColor(255, 255, 255)     # Primary headers
    TEXT_MUTED = RGBColor(165, 180, 206)     # Secondary descriptions
    
    BLUE_BRAND = RGBColor(59, 130, 246)      # Triage Blue
    CYAN_BRAND = RGBColor(15, 185, 196)      # Hypothesis Cyan
    GREEN_BRAND = RGBColor(34, 197, 94)      # Evidence Green
    VIOLET_BRAND = RGBColor(139, 92, 246)    # Critic Violet
    AMBER_BRAND = RGBColor(245, 158, 11)     # Judge Amber
    RED_ERROR = RGBColor(239, 68, 68)        # Problem Red
    
    # Helper to set slide background
    def set_slide_background(slide):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = BG_COLOR

    # Helper to create a premium card
    def add_card(slide, left, top, width, height, bg_rgb=CARD_BG, border_rgb=BORDER_COLOR):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_rgb
        shape.line.color.rgb = border_rgb
        shape.line.width = Pt(1.5)
        
        # Text Frame config
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.2)
        tf.margin_bottom = Inches(0.2)
        return shape

    # Helper to add standard header to a slide
    def add_slide_header(slide, title_text, category="SENTRYSWARM AI"):
        set_slide_background(slide)
        
        # Add textbox for category and title
        header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(1.0))
        tf = header_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        
        p_cat = tf.paragraphs[0]
        p_cat.text = category.upper()
        p_cat.font.name = 'Segoe UI'
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = BLUE_BRAND
        p_cat.space_after = Pt(2)
        
        p_title = tf.add_paragraph()
        p_title.text = title_text
        p_title.font.name = 'Segoe UI'
        p_title.font.size = Pt(26)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_WHITE

    # ----------------------------------------------------
    # SLIDE 1: Hero Title Slide (Visual Outage Swarm concept)
    # ----------------------------------------------------
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide1)
    
    # Left Hero Text
    hero_box = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(7.5), Inches(4.5))
    tf1 = hero_box.text_frame
    tf1.word_wrap = True
    tf1.margin_left = tf1.margin_right = tf1.margin_top = tf1.margin_bottom = 0
    
    p_tag = tf1.paragraphs[0]
    p_tag.text = "MICROSOFT BUILD AI HACKATHON  ·  THEME: AGENT SWARMS"
    p_tag.font.name = 'Segoe UI'
    p_tag.font.size = Pt(12)
    p_tag.font.bold = True
    p_tag.font.color.rgb = CYAN_BRAND
    p_tag.space_after = Pt(12)
    
    p_main = tf1.add_paragraph()
    p_main.text = "SentrySwarm AI"
    p_main.font.name = 'Segoe UI'
    p_main.font.size = Pt(64)
    p_main.font.bold = True
    p_main.font.color.rgb = TEXT_WHITE
    p_main.space_after = Pt(4)
    
    p_sub = tf1.add_paragraph()
    p_sub.text = "Autonomous Outage Diagnosis & Resolution"
    p_sub.font.name = 'Segoe UI'
    p_sub.font.size = Pt(22)
    p_sub.font.color.rgb = BLUE_BRAND
    p_sub.space_after = Pt(24)
    
    p_desc = tf1.add_paragraph()
    p_desc.text = "An orchestrator dispatching specialized AI agents that collaborate and debate in parallel to find outage root causes—live streamed in real time."
    p_desc.font.name = 'Segoe UI'
    p_desc.font.size = Pt(15)
    p_desc.font.color.rgb = TEXT_MUTED
    p_desc.space_after = Pt(36)
    
    p_devs = tf1.add_paragraph()
    p_devs.text = "Developed by: Avinash Yadav & Darshan Ghorpade"
    p_devs.font.name = 'Segoe UI'
    p_devs.font.size = Pt(16)
    p_devs.font.bold = True
    p_devs.font.color.rgb = TEXT_WHITE

    # Right visual accent container
    viz_card = add_card(slide1, Inches(8.8), Inches(1.8), Inches(3.7), Inches(4.5))
    tf_v = viz_card.text_frame
    tf_v.margin_top = Inches(0.5)
    
    p_vt = tf_v.paragraphs[0]
    p_vt.text = "🧭 Triage Agent\n      │\n┌─────┴─────┐\n▼     ▼     ▼\n💡    💡    💡 (Hypotheses)\n│     │     │\n▼     ▼     ▼\n🔍    🔍    🔍 (Evidence)\n└─────┬─────┘\n      ▼\n⚔️ Critic Agent\n      │\n      ▼\n⚖️ Judge Verdict"
    p_vt.font.name = 'Consolas'
    p_vt.font.size = Pt(13)
    p_vt.font.bold = True
    p_vt.font.color.rgb = CYAN_BRAND
    p_vt.alignment = PP_ALIGN.CENTER

    # ----------------------------------------------------
    # SLIDE 2: The Problem
    # ----------------------------------------------------
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide2, "Outage Diagnosis is Costly and Slow")
    
    # Left Card: The Outage Cost (Red Brand)
    card_l = add_card(slide2, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0), border_rgb=RED_ERROR)
    tf_l = card_l.text_frame
    
    p_lt = tf_l.paragraphs[0]
    p_lt.text = "🚨  The Crisis: High MTTR"
    p_lt.font.name = 'Segoe UI'
    p_lt.font.size = Pt(20)
    p_lt.font.bold = True
    p_lt.font.color.rgb = RED_ERROR
    p_lt.space_after = Pt(18)
    
    bullets_l = [
        "System outages cost modern enterprises an average of $9,000 per minute.",
        "Manual war-rooms suffer from cognitive overload, trying to read thousands of disconnected log entries and metrics simultaneously.",
        "Incident responders rely on intuition and confirmation bias, leading to delayed remediation and prolonged downtime."
    ]
    for b in bullets_l:
        p_b = tf_l.add_paragraph()
        p_b.text = "•  " + b
        p_b.font.name = 'Segoe UI'
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = TEXT_MUTED
        p_b.space_before = Pt(12)
        
    # Right Card: The Single AI Agent Limitation
    card_r = add_card(slide2, Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.0))
    tf_r = card_r.text_frame
    
    p_rt = tf_r.paragraphs[0]
    p_rt.text = "⚠️  The Limitation of Single AI Agents"
    p_rt.font.name = 'Segoe UI'
    p_rt.font.size = Pt(20)
    p_rt.font.bold = True
    p_rt.font.color.rgb = AMBER_BRAND
    p_rt.space_after = Pt(18)
    
    bullets_r = [
        "A single AI agent with tools tends to accept its first assumption without validation.",
        "Single-agent loops frequently get stuck, generating repeating tool-calling commands when faced with noise.",
        "Without an adversarial counterpart, single LLM reasoning cannot guarantee defensible, hallucination-free answers."
    ]
    for b in bullets_r:
        p_b = tf_r.add_paragraph()
        p_b.text = "•  " + b
        p_b.font.name = 'Segoe UI'
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = TEXT_MUTED
        p_b.space_before = Pt(12)

    # ----------------------------------------------------
    # SLIDE 3: Swarm Pipeline Architecture
    # ----------------------------------------------------
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide3, "The Collaborative SentrySwarm Pipeline")
    
    # 5 cards aligned horizontally, each represents an agent.
    roles = [
        {"icon": "🧭", "name": "Triage", "color": BLUE_BRAND, "desc": "Analyzes the initial alert symptoms, logs, and classifies severity. Spawns investigators."},
        {"icon": "💡", "name": "Hypothesis", "color": CYAN_BRAND, "desc": "N investigators run in parallel, formulating distinct root-cause possibilities from separate angles."},
        {"icon": "🔍", "name": "Evidence", "color": GREEN_BRAND, "desc": "Fact-checks hypotheses by executing regex searches over production logs, separating signal from noise."},
        {"icon": "⚔️", "name": "Critic", "color": VIOLET_BRAND, "desc": "Red-teams hypotheses, exposing cognitive leaps, unverified claims, or unproven assumptions."},
        {"icon": "⚖️", "name": "Judge", "color": AMBER_BRAND, "desc": "Weighs hypotheses against evidence and critiques. Converges on verdict and remediation script."}
    ]
    
    card_w = Inches(2.2)
    card_h = Inches(4.7)
    gap = Inches(0.18)
    start_x = Inches(0.8)
    top_y = Inches(1.8)
    
    for idx, r in enumerate(roles):
        left_x = start_x + idx * (card_w + gap)
        # Create card with customized border matching the role
        c_shape = add_card(slide3, left_x, top_y, card_w, card_h, border_rgb=r["color"])
        tf_c = c_shape.text_frame
        
        p_icon = tf_c.paragraphs[0]
        p_icon.text = r["icon"]
        p_icon.font.name = 'Segoe UI Symbol'
        p_icon.font.size = Pt(40)
        p_icon.alignment = PP_ALIGN.CENTER
        
        p_name = tf_c.add_paragraph()
        p_name.text = r["name"]
        p_name.font.name = 'Segoe UI'
        p_name.font.size = Pt(18)
        p_name.font.bold = True
        p_name.font.color.rgb = r["color"]
        p_name.alignment = PP_ALIGN.CENTER
        p_name.space_before = Pt(8)
        p_name.space_after = Pt(12)
        
        p_desc = tf_c.add_paragraph()
        p_desc.text = r["desc"]
        p_desc.font.name = 'Segoe UI'
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = TEXT_MUTED
        p_desc.alignment = PP_ALIGN.LEFT
        p_desc.space_before = Pt(6)
        
        # Add visual connector arrow between cards (except last)
        if idx < 4:
            arrow_box = slide3.shapes.add_textbox(left_x + card_w, top_y + Inches(2.0), gap, Inches(0.8))
            tf_a = arrow_box.text_frame
            tf_a.margin_left = tf_a.margin_right = tf_a.margin_top = tf_a.margin_bottom = 0
            p_a = tf_a.paragraphs[0]
            p_a.text = "→"
            p_a.font.name = 'Segoe UI'
            p_a.font.size = Pt(22)
            p_a.font.bold = True
            p_a.font.color.rgb = BORDER_COLOR
            p_a.alignment = PP_ALIGN.CENTER

    # ----------------------------------------------------
    # SLIDE 4: Technical core / Innovations
    # ----------------------------------------------------
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide4, "Core Technical Innovations")
    
    innovations = [
        {"title": "Real-Time WebSocket Graph Rendering", "desc": "A customized vanilla HTML5 SVG graph engine aligns connections dynamically using requestAnimationFrame. streams agent output token-by-token directly to the canvas so developers watch the swarm think live.", "color": BLUE_BRAND},
        {"title": "Interactive Out-of-Band Rerun Retries", "desc": "A custom WebSocket payload handler captures user requests to retry specific nodes. It compiles current hypotheses and evidence states into a context bundle and runs individual agent tasks on the fly.", "color": GREEN_BRAND},
        {"title": "Multi-Model Context Propagation", "desc": "Leverages Python ContextVars to dynamically route calls across Gemini 3.5, Gemma 4, and Anthropic Claude models on a single execution thread without mutating core agent logic.", "color": CYAN_BRAND},
        {"title": "Zero-Config Client-Side Privacy", "desc": "Keeps API keys out of backend database storage. Keys are configured locally, stored in the browser's localStorage, and securely injected during WebSocket connections.", "color": VIOLET_BRAND}
    ]
    
    for idx, inn in enumerate(innovations):
        row = idx // 2
        col = idx % 2
        
        left = Inches(0.8 + col * 6.0)
        top = Inches(1.7 + row * 2.6)
        
        card = add_card(slide4, left, top, Inches(5.7), Inches(2.3))
        tf_i = card.text_frame
        
        p_t = tf_i.paragraphs[0]
        p_t.text = "⚡  " + inn["title"]
        p_t.font.name = 'Segoe UI'
        p_t.font.size = Pt(18)
        p_t.font.bold = True
        p_t.font.color.rgb = inn["color"]
        p_t.space_after = Pt(8)
        
        p_d = tf_i.add_paragraph()
        p_d.text = inn["desc"]
        p_d.font.name = 'Segoe UI'
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = TEXT_MUTED
        p_d.space_before = Pt(4)

    # ----------------------------------------------------
    # SLIDE 5: Quick Start / Outage Demo Case
    # ----------------------------------------------------
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide5, "Execution, Live Demo & Setup")
    
    # Left Card: Live Demo Case
    card_demo = add_card(slide5, Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.0))
    tf_d = card_demo.text_frame
    
    p_dt = tf_d.paragraphs[0]
    p_dt.text = "🎯  Built-in Incident: INC-4821"
    p_dt.font.name = 'Segoe UI'
    p_dt.font.size = Pt(20)
    p_dt.font.bold = True
    p_dt.font.color.rgb = GREEN_BRAND
    p_dt.space_after = Pt(14)
    
    p_dd = tf_d.add_paragraph()
    p_dd.text = "Includes realistic test cases to demonstrate Swarm reasoning:"
    p_dd.font.name = 'Segoe UI'
    p_dd.font.size = Pt(13)
    p_dd.font.color.rgb = TEXT_WHITE
    p_dd.space_after = Pt(8)
    
    points = [
        "Outage Symptom: Latency spike and elevated 500 errors on checkout service.",
        "Root Cause: A database deployment changes max connections from 50 to 10.",
        "Red Herrings: Placed fake Redis key-eviction logs and payment gateway gateway-timeout warnings.",
        "Outcome: SentrySwarm's Critic successfully rejects the red herrings, and the Judge pinpoints connection-pool exhaustion."
    ]
    for p in points:
        p_pt = tf_d.add_paragraph()
        p_pt.text = "✔  " + p
        p_pt.font.name = 'Segoe UI'
        p_pt.font.size = Pt(12.5)
        p_pt.font.color.rgb = TEXT_MUTED
        p_pt.space_before = Pt(8)

    # Right Card: Run commands
    card_run = add_card(slide5, Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.0))
    tf_ru = card_run.text_frame
    
    p_rut = tf_ru.paragraphs[0]
    p_rut.text = "⚙  Deploy & Run SentrySwarm"
    p_rut.font.name = 'Segoe UI'
    p_rut.font.size = Pt(20)
    p_rut.font.bold = True
    p_rut.font.color.rgb = BLUE_BRAND
    p_rut.space_after = Pt(14)
    
    run_steps = [
        "1. Clone the project code:\n   git clone https://github.com/AvizRepo/Sentry-Swarm-Agent.git",
        "2. Run locally via Python:\n   pip install -r requirements.txt\n   cd backend && uvicorn main:app --reload",
        "3. Run via Docker Compose:\n   docker compose up --build",
        "4. Setup Keys & Diagnose:\n   Open http://localhost:8000, add API keys in '🔑 Key' modal, and click 'Run Swarm'."
    ]
    for step in run_steps:
        p_step = tf_ru.add_paragraph()
        p_step.text = step
        p_step.font.name = 'Segoe UI'
        p_step.font.size = Pt(12)
        p_step.font.color.rgb = TEXT_MUTED
        p_step.space_before = Pt(10)

    # Save the presentation
    prs.save("sentryswarm_pitch.pptx")
    print("Beautiful sentryswarm_pitch.pptx created successfully!")

if __name__ == '__main__':
    create_presentation()
